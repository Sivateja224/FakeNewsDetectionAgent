# generate_presentation.py
# This script programmatically generates a professional, college final-year level
# PowerPoint presentation (.pptx) summarizing the Fake News Detection Agent project.
# It uses the python-pptx library to construct the slide deck.

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    print("=== Creating PowerPoint Presentation ===")
    prs = Presentation()
    
    # Define a professional color palette matching the application
    PRIMARY_BLUE = RGBColor(10, 58, 96)     # Deep Blue
    SECONDARY_BLUE = RGBColor(41, 128, 185) # Electric Blue
    DARK_TEXT = RGBColor(44, 62, 80)        # Dark Charcoal
    LIGHT_TEXT = RGBColor(127, 140, 141)    # Muted Gray
    WHITE = RGBColor(255, 255, 255)
    
    # Helper function to set slide background to a clean layout
    def set_slide_title(slide, text, color=PRIMARY_BLUE):
        title_shape = slide.shapes.title
        title_shape.text = text
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.name = 'Arial'
            paragraph.font.size = Pt(36)
            paragraph.font.bold = True
            paragraph.font.color.rgb = color
            
    # Helper to add standard bullet points
    def add_bullets(slide, points, size=18, spacing=10):
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        for i, point_text in enumerate(points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point_text
            p.font.name = 'Arial'
            p.font.size = Pt(size)
            p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(spacing)

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Layout 0)
    # ----------------------------------------------------
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = "Fake News Detection Agent"
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    subtitle_shape.text = (
        "An Agentic AI Web Application Using Machine Learning\n"
        "College Final-Year Project Submission\n"
        "Powered by Flask, Scikit-Learn, and Bootstrap 5"
    )
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(18)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = SECONDARY_BLUE
    
    # ----------------------------------------------------
    # SLIDE 2: Project Overview & Objectives (Layout 1)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "1. Project Overview & Objectives")
    bullets_2 = [
        "Objective: Develop an intelligent Agentic AI web application that evaluates news articles in real-time.",
        "Problem: Misinformation and fake news spread exponentially online, requiring automated detection mechanisms.",
        "Agentic Behavior: Unlike traditional classifiers, this system not only labels news (Real/Fake) but also computes confidence, extracts linguistic indicators, and provides recommendations.",
        "Target Audience: General users, journalists, and students who require quick validation of internet claims."
    ]
    add_bullets(slide, bullets_2)

    # ----------------------------------------------------
    # SLIDE 3: System Architecture (Layout 1)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "2. System Architecture")
    bullets_3 = [
        "Model Type: Binary Classification via Passive Aggressive Classifier.",
        "Feature Representation: Word-level TF-IDF Vectorization with unigrams and bigrams (5,000 feature limit).",
        "Backend Architecture: Flask-based micro-framework running Python 3.11.x.",
        "Frontend Interface: Clean single-page application using Bootstrap 5 with responsive grid design.",
        "Failsafe Mechanism: Local dataset fallback, and pre-trained pickle/joblib modules loaded at server startup."
    ]
    add_bullets(slide, bullets_3)

    # ----------------------------------------------------
    # SLIDE 4: Machine Learning Pipeline (Layout 1)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "3. Machine Learning Pipeline")
    bullets_4 = [
        "Data Merging: Combined Fake.csv (0 label) and True.csv (1 label) into a unified dataframe.",
        "Linguistic Cleaning: Lowercased all text, removed punctuation, digits, and standard NLTK stopwords.",
        "Lemmatization: Applied WordNet Lemmatizer to scale words down to base lexical forms.",
        "Vectorization: Built an TF-IDF feature matrix based on text frequency inverse document frequency.",
        "Split & Train: Trained Passive Aggressive Classifier on 80% split with 20% test partition for evaluation."
    ]
    add_bullets(slide, bullets_4)

    # ----------------------------------------------------
    # SLIDE 5: Dataset Information (Layout 1)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "4. Dataset Information")
    bullets_5 = [
        "Dataset Sources: Aggregated set of real and fake articles covering news across politics, world news, and local affairs.",
        "True.csv size: 53.5 MB (21,417 articles with title, text, subject, date).",
        "Fake.csv size: 62.7 MB (23,481 articles with title, text, subject, date).",
        "Combined Total: 44,898 records.",
        "Feature Combination: Article 'title' and 'body text' are merged to improve classification signals."
    ]
    add_bullets(slide, bullets_5)

    # ----------------------------------------------------
    # SLIDE 6: Agent Workflow (Layout 1)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "5. Agent Workflow")
    bullets_6 = [
        "User Input: The user pastes an article title/content into the web dashboard.",
        "AI Analysis: Preprocesses the text (lowercasing, cleaning, lemmatizing) and applies TF-IDF weights.",
        "Inference: Generates classification prediction (Real/Fake).",
        "Confidence Estimation: Converts decision function distance to a percentage using a Sigmoid function.",
        "Explanation & Flags: Highlights suspicious text markers (exclamation, shout-caps, clickbait keys).",
        "Actionable Advice: Recommends verifying claims using reputable sources."
    ]
    add_bullets(slide, bullets_6)

    # ----------------------------------------------------
    # SLIDE 7: Web Application Features (Layout 1)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "6. Web Application Features")
    bullets_7 = [
        "Glassmorphism UI: Trendy visual appearance using blur backdrops and clean rounded cards.",
        "Responsive Grid: Mobile-friendly view optimized for tablets, laptops, and smart screens.",
        "Interactive Status Indicators: Visual progress bar and color-coded alert panels.",
        "Prediction History: Current session log tracking inputs and predictions on the dashboard.",
        "Clear State Controls: Instantly reset input text boxes, results cards, and confidence meters."
    ]
    add_bullets(slide, bullets_7)

    # ----------------------------------------------------
    # SLIDE 8: Future Scope (Layout 1)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "7. Future Scope & Enhancements")
    bullets_8 = [
        "Deep Learning: Upgrade model to Transformer-based architectures (BERT, RoBERTa) for contextual semantic understanding.",
        "API Integrations: Connect with live web-scraping search APIs to check primary sources in real time.",
        "Browser Extensions: Port model inference to chrome extension format to classify text directly on social platforms.",
        "Multilingual Support: Enable translation and classification of non-English news articles."
    ]
    add_bullets(slide, bullets_8)

    # ----------------------------------------------------
    # SLIDE 9: Conclusion (Layout 1)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "8. Conclusion")
    bullets_9 = [
        "Summary: Successfully built a modern Fake News Detection Agent that bridges Machine Learning classification with Agentic AI output explanations.",
        "Technology Validation: Validated using high precision, recall, and accuracy metrics (~99% in PAC evaluations).",
        "College Submission: Designed with professional-grade software structures, documentation, and web design, making it fully ready for final-year submission."
    ]
    add_bullets(slide, bullets_9)

    # Save presentation
    output_filename = "presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as '{output_filename}'!")

if __name__ == "__main__":
    create_presentation()
